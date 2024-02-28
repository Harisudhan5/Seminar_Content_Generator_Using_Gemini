from pptx import Presentation

pr1 = Presentation()

# Slide 1: Title, Introduction Line
slide_1_layout = pr1.slide_layouts[0]
slide_1 = pr1.slides.add_slide(slide_1_layout)
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Presentation Using Python Code"
subtitle_1.text = "Introduction to the Course"

# Slide 2: Title, 7 Bullet Points
slide_2_layout = pr1.slide_layouts[1]
slide_2 = pr1.slides.add_slide(slide_2_layout)
title_2 = slide_2.shapes.title
title_2.text = "Seven Bullet Points"
bullet_point_box = slide_2.placeholders[1]

bullet_points = [
    "Bullet Point 1",
    "Bullet Point 2",
    "Bullet Point 3",
    "Bullet Point 4",
    "Bullet Point 5",
    "Bullet Point 6",
    "Bullet Point 7",
]

for point in bullet_points:
    bullet_point_box.text += f"\n- {point}"

# Slide 3: Title, Paragraph
slide_3_layout = pr1.slide_layouts[1]  # You can choose a different layout if needed
slide_3 = pr1.slides.add_slide(slide_3_layout)
title_3 = slide_3.shapes.title
title_3.text = "Slide with Paragraph"
paragraph_box = slide_3.placeholders[1]
paragraph_box.text = "This is a paragraph of text. You can add more details and information here."

# Save the presentation
pr1.save("Updated_Presentation.pptx")
